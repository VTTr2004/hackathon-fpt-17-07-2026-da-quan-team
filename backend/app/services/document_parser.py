from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md", ".csv", ".json"}:
        return path.read_text(encoding="utf-8", errors="replace")
    if suffix == ".pdf":
        reader = PdfReader(str(path))
        return "\n\n".join(f"[PAGE {index}]\n{page.extract_text() or ''}" for index, page in enumerate(reader.pages, 1))
    if suffix == ".docx":
        document = DocxDocument(str(path))
        sections = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for index, table in enumerate(document.tables, 1):
            rows = ["\t".join(cell.text.strip() for cell in row.cells) for row in table.rows]
            sections.append(f"[TABLE {index}]\n" + "\n".join(rows))
        return "\n\n".join(sections)
    if suffix == ".pptx":
        presentation = Presentation(str(path))
        slides = []
        for index, slide in enumerate(presentation.slides, 1):
            parts: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    parts.extend("\t".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows)
                elif hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
            text = "\n".join(parts)
            slides.append(f"[SLIDE {index}]\n{text}")
        return "\n\n".join(slides)
    if suffix == ".xlsx":
        workbook = load_workbook(path, read_only=True, data_only=True)
        sheets = []
        for sheet in workbook.worksheets:
            rows = ["\t".join("" if value is None else str(value) for value in row) for row in sheet.values]
            sheets.append(f"[SHEET {sheet.title}]\n" + "\n".join(rows))
        return "\n\n".join(sheets)
    raise ValueError(f"Định dạng {suffix or 'không xác định'} chưa được hỗ trợ")
